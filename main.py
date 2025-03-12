from kivy.app import App
from kivy.uix.boxlayout import BoxLayout
from kivy.uix.label import Label
from kivy.uix.button import Button
from kivy.uix.popup import Popup
from kivy.uix.image import AsyncImage
from kivy.graphics import Color, Rectangle, RoundedRectangle
from kivy.metrics import dp, sp
from kivy.clock import Clock
from kivy.uix.boxlayout import BoxLayout
from kivy.uix.scrollview import ScrollView
from kivy.core.window import Window
import os
import tempfile
from threading import Thread
from pptx import Presentation
from PIL import Image, ImageDraw, ImageFont
import io
import fitz
from kivy.config import Config
Config.set('kivy', 'window_icon', '')
Config.set('graphics', 'resizable', '0')
Config.set('graphics', 'width', '412')
Config.set('graphics', 'height', '732')
Config.set('graphics', 'orientation', 'portrait')

Window.rotation = 0
Window.size = (412, 615)

class MainMenu(BoxLayout):
    def __init__(self, **kwargs):
        super().__init__(orientation='vertical', **kwargs)
        self.loading_popup = None
        self.current_navigation_path = ""  # Track current navigation path

        # Orange background
        with self.canvas.before:
            Color(1, 0.6, 0.2, 1)  # Orange
            self.bg_rect = Rectangle(size=self.size, pos=self.pos)
        self.bind(size=self._update_bg, pos=self._update_bg)

        # Top spacer
        self.add_widget(BoxLayout(size_hint=(1, 0.2)))

        # White card - make size_hint responsive instead of fixed width
        self.card = BoxLayout(
            orientation='vertical',
            padding=dp(20),
            spacing=dp(10),
            size_hint=(0.9, None),
            pos_hint={'center_x': 0.5}
        )
        self.card.bind(minimum_height=self.card.setter('height'))

        with self.card.canvas.before:
            Color(1, 1, 1, 1)  # White
            self.card.rect = RoundedRectangle(
                size=self.card.size,
                pos=self.card.pos,
                radius=[dp(15)]
            )
        self.card.bind(size=self._update_card_bg, pos=self._update_card_bg)

        # Title - use adaptive sizing
        title_label = Label(
            text="Gregor-E Library",
            font_size=self.get_adaptive_font_size(28),
            bold=True,
            color=(0, 0, 0, 1),
            size_hint=(1, None),
            height=dp(40)
        )
        self.card.add_widget(title_label)

        # Spacer
        self.card.add_widget(BoxLayout(size_hint=(1, None), height=dp(5)))

        # Tagline - use adaptive sizing
        subheading_label = Label(
            text="Your digital library at your fingertips",
            font_size=self.get_adaptive_font_size(14),
            color=(0, 0, 0, 1),
            size_hint=(1, None),
            height=dp(30)
        )
        self.card.add_widget(subheading_label)

        # Spacer
        self.card.add_widget(BoxLayout(size_hint=(1, None), height=dp(10)))

        # Buttons
        self._add_pill_button("Grade 11", self.show_grade11)
        self._add_pill_button("Grade 12", self.show_grade12)
        self._add_pill_button("About Us", self.show_about_us)

        # Add card to main layout
        self.add_widget(self.card)

        # Bottom spacer
        self.add_widget(BoxLayout(size_hint=(1, 0.2)))

        # Update layout when window size changes
        Window.bind(on_resize=self.on_window_resize)

    def on_window_resize(self, instance, width, height):
        """Handle window resize events to update adaptive layouts"""
        # Update card size based on window dimensions
        self.card.height = self.card.minimum_height

        # Update fonts
        for child in self.card.children:
            if isinstance(child, Label):
                if "Gregor-E Library" in child.text:
                    child.font_size = self.get_adaptive_font_size(28)
                else:
                    child.font_size = self.get_adaptive_font_size(14)
            elif isinstance(child, Button):
                child.font_size = self.get_adaptive_font_size(16)

    def get_adaptive_font_size(self, base_size):
        """Calculate font size based on screen width"""
        # Get the smaller dimension (portrait vs landscape)
        min_dimension = min(Window.width, Window.height)
        # Base calculation for adaptive font size
        # Reference width: 360dp (common phone width)
        reference_width = 360
        return sp(base_size * (min_dimension / reference_width) * 0.7)

    def _update_bg(self, *args):
        self.bg_rect.pos = self.pos
        self.bg_rect.size = self.size

    def _update_card_bg(self, instance, *args):
        instance.rect.pos = instance.pos
        instance.rect.size = instance.size

    def _add_pill_button(self, text, callback):
        btn = Button(
            text=text,
            size_hint=(1, None),
            height=dp(50),  # Increased touch target for mobile
            font_size=self.get_adaptive_font_size(16),
            color=(0, 0, 0, 1),
            background_color=(0, 0, 0, 0)
        )
        btn.bind(on_release=callback)

        with btn.canvas.before:
            Color(0.8, 0.8, 0.8, 1)  # Light gray
            btn.rect = RoundedRectangle(
                size=btn.size,
                pos=btn.pos,
                radius=[dp(25)]  # Larger radius for better look on mobile
            )
        btn.bind(size=self._update_btn_shape, pos=self._update_btn_shape)

        self.card.add_widget(btn)

    def _update_btn_shape(self, btn, *args):
        btn.rect.pos = btn.pos
        btn.rect.size = btn.size

    def show_grade11(self, instance):
        self._show_quarters_popup("Grade 11")

    def show_grade12(self, instance):
        self._show_quarters_popup("Grade 12")

    def _show_quarters_popup(self, title):
        self.current_navigation_path = title

        popup_content = BoxLayout(orientation='vertical', padding=dp(20), spacing=dp(15))
        self._set_dark_bg(popup_content)

        popup_content.add_widget(Label(
            text=f"{title} - Select Quarter",
            font_size=self.get_adaptive_font_size(20),
            color=(1, 1, 1, 1),
            size_hint=(1, None),
            height=dp(40)
        ))

        quarters = ["1st", "2nd", "3rd", "4th"]
        for i, quarter in enumerate(quarters, start=1):
            quarter_button = Button(
                text=f"📖 {quarter} Quarter",
                size_hint=(1, None),
                height=dp(50),  # Increased height for better touch targets
                font_size=self.get_adaptive_font_size(16)
            )
            # Fix: Capture the loop variable i in the lambda function
            quarter_num = i  # Explicitly get the quarter number
            quarter_button.bind(on_release=lambda btn, q=quarter_num:
            self._show_subjects_popup(f"{title} - {quarters[q - 1]} Quarter", q))
            self._style_pill_button(quarter_button, (0.5, 0.5, 0.5, 1), (1, 1, 1, 1))
            popup_content.add_widget(quarter_button)

        close_button = Button(
            text="❌ Close",
            size_hint=(1, None),
            height=dp(50),
            font_size=self.get_adaptive_font_size(16)
        )
        self._style_pill_button(close_button, (0.8, 0.2, 0.2, 1), (1, 1, 1, 1))

        # Make popup responsive with size_hint
        popup = Popup(
            title=title,
            content=popup_content,
            size_hint=(0.95, 0.8),
            title_size=self.get_adaptive_font_size(18)
        )
        close_button.bind(on_release=popup.dismiss)
        popup_content.add_widget(close_button)
        popup.open()

    def _show_subjects_popup(self, title, quarter_num=None):
        # Store current navigation path for later use
        self.current_navigation_path = title

        # Extract quarter from title if not provided
        if quarter_num is None:
            # Gets the '1' from '1st'
            quarter_part = title.split(" - ")[1].split()[0][0]
            quarter_num = int(quarter_part)

        popup_content = BoxLayout(orientation='vertical', padding=dp(20), spacing=dp(15))
        self._set_dark_bg(popup_content)

        popup_content.add_widget(Label(
            text=f"{title} - Select Subject",
            font_size=self.get_adaptive_font_size(14),
            color=(1, 1, 1, 1),
            size_hint=(1, None),
            height=dp(40)
        ))

        # Determine grade from title
        grade = 11 if "Grade 11" in title else 12

        # Get subjects based on grade and quarter
        if grade == 11:
            if quarter_num in (1, 2):
                subjects = [
                    "Oral Comm",
                    "Komunikasyon",
                    "General Mathematics",
                    "Earth Science",
                    "Understanding Culture, Society and Politics",
                    "Personal Development",
                    "PE"
                ]
            else:
                subjects = [
                    "Reading and Writing",
                    "Pagbasa",
                    "21st Century Literature",
                    "Stats and Probability",
                    "PE"
                ]
        else:  # Grade 12
            subjects = [
                "1 - Introduction & Overview",
                "2 - Key Concepts & Theories",
                "3 - Practical Applications",
                "4 - Review and Assessment",
                "5 - Advanced Topics or Follow-Up",
                "6 - Additional Case Studies",
                "7 - Problem-Solving Approaches",
                "8 - Recap and Exam Preparation"
            ]

            subjects = subjects[:5] if quarter_num in (1, 2) else subjects[5:]

        # Create scroll view for subjects
        scroll_height = min(Window.height * 0.5, dp(300))
        scroll_view = ScrollView(size_hint=(1, None), height=scroll_height)
        subjects_layout = BoxLayout(orientation='vertical', size_hint_y=None, spacing=dp(15))
        subjects_layout.bind(minimum_height=subjects_layout.setter('height'))

        for subject in subjects:
            subject_button = Button(
                text=f"📚 {subject}",
                size_hint=(1, None),
                height=dp(50),
                font_size=self.get_adaptive_font_size(12)
            )
            # Capture current subject in lambda
            subject_button.bind(on_release=lambda btn, q=quarter_num, s=subject:
            self._show_weeks_popup(f"{title} - {s}", q, s))
            self._style_pill_button(subject_button, (0.6, 0.6, 0.6, 1), (1, 1, 1, 1))
            subjects_layout.add_widget(subject_button)

        scroll_view.add_widget(subjects_layout)
        popup_content.add_widget(scroll_view)

        # Close button and popup creation remains the same
        close_button = Button(
            text="❌ Close",
            size_hint=(1, None),
            height=dp(50),
            font_size=self.get_adaptive_font_size(12)
        )
        self._style_pill_button(close_button, (0.8, 0.2, 0.2, 1), (1, 1, 1, 1))

        popup = Popup(
            title=title,
            content=popup_content,
            size_hint=(0.95, 0.9),
            title_size=self.get_adaptive_font_size(14)
        )
        close_button.bind(on_release=popup.dismiss)
        popup_content.add_widget(close_button)
        popup.open()

    def _show_weeks_popup(self, title, quarter_num=None, subject=None):
        self.current_navigation_path = title

        # Extract quarter and subject if not provided
        if quarter_num is None:
            title_parts = title.split(" - ")
            quarter_part = title_parts[1].split()[0][0]  # Gets the '1' from '1st'
            quarter_num = int(quarter_part)

        if subject is None:
            subject = title.split(" - ")[-1]

        popup_content = BoxLayout(orientation='vertical', padding=dp(20), spacing=dp(15))
        self._set_dark_bg(popup_content)

        popup_content.add_widget(Label(
            text=f"{title} - Select Week",
            font_size=self.get_adaptive_font_size(14),
            color=(1, 1, 1, 1),
            size_hint=(1, None),
            height=dp(40)
        ))

        # Create scrollable area for weeks
        scroll_height = min(Window.height * 0.6, dp(400))  # Adaptive height
        scroll_view = ScrollView(size_hint=(1, None), height=scroll_height)
        weeks_layout = BoxLayout(orientation='vertical', size_hint_y=None, spacing=dp(15))
        weeks_layout.bind(minimum_height=weeks_layout.setter('height'))

        # Use 8 weeks to match the batch script
        weeks = [f"Week {i}" for i in range(1, 9)]
        for week in weeks:
            week_button = Button(
                text=f"📅 {week}",
                size_hint=(1, None),
                height=dp(50),
                font_size=self.get_adaptive_font_size(14)
            )
            # Fix: Properly capture week number in the lambda
            week_num = int(week.split()[1])
            week_button.bind(on_release=lambda btn, q=quarter_num, s=subject, w=week_num:
            self._show_files_popup(f"{title} - Week {w}", q, s, w))
            self._style_pill_button(week_button, (0.7, 0.7, 0.7, 1), (1, 1, 1, 1))
            weeks_layout.add_widget(week_button)

        scroll_view.add_widget(weeks_layout)
        popup_content.add_widget(scroll_view)

        close_button = Button(
            text="❌ Close",
            size_hint=(1, None),
            height=dp(50),
            font_size=self.get_adaptive_font_size(14)
        )
        self._style_pill_button(close_button, (0.8, 0.2, 0.2, 1), (1, 1, 1, 1))

        # Responsive popup
        popup = Popup(
            title=title,
            content=popup_content,
            size_hint=(0.95, 0.9),
            title_size=self.get_adaptive_font_size(14)
        )
        close_button.bind(on_release=popup.dismiss)
        popup_content.add_widget(close_button)
        popup.open()

    def _show_files_popup(self, title, quarter_num=None, subject=None, week_num=None):
        self.current_navigation_path = title

        # Extract information if not provided
        if any(param is None for param in [quarter_num, subject, week_num]):
            title_parts = title.split(" - ")
            if quarter_num is None:
                quarter_part = title_parts[1].split()[0][0]  # Gets the '1' from '1st'
                quarter_num = int(quarter_part)

            if subject is None:
                subject = title_parts[2].replace("Subject ", "")

            if week_num is None:
                week_num = int(title_parts[3].split()[1])

        popup_content = BoxLayout(orientation='vertical', padding=dp(20), spacing=dp(15))
        self._set_dark_bg(popup_content)

        popup_content.add_widget(Label(
            text=f"{title} - Select File",
            font_size=self.get_adaptive_font_size(20),
            color=(1, 1, 1, 1),
            size_hint=(1, None),
            height=dp(40)
        ))

        # Create scrollable area for files
        scroll_height = min(Window.height * 0.6, dp(400))  # Adaptive height
        scroll_view = ScrollView(size_hint=(1, None), height=scroll_height)
        files_layout = BoxLayout(orientation='vertical', size_hint_y=None, spacing=dp(15))
        files_layout.bind(minimum_height=files_layout.setter('height'))

        BASE_DIR = "GregorELibrary"

        folder_path = os.path.join(
            BASE_DIR,
            f"Quarter {quarter_num}",
            f"Subject {subject}",
            f"Week {week_num}"
        )

        # Debug output to verify path components
        print(f"Looking for files in: {folder_path}")
        print(f"Quarter: {quarter_num}, Subject: {subject}, Week: {week_num}")

        files = []
        if os.path.exists(folder_path):
            # Include both PPTX and PDF files
            files = [f for f in os.listdir(folder_path)
                     if f.lower().endswith(('.pptx', '.pdf'))]

        if files:
            for file in files:
                file_button = Button(
                    text=f"📄 {file}",
                    size_hint=(1, None),
                    height=dp(50),
                    font_size=self.get_adaptive_font_size(14)
                )
                # Fix: Properly capture file name in the lambda
                file_button.bind(on_release=lambda btn, f=file, q=quarter_num, s=subject, w=week_num:
                self._handle_file_selected(f, q, s, w))
                self._style_pill_button(file_button, (0.8, 0.8, 0.8, 1), (1, 1, 1, 1))
                files_layout.add_widget(file_button)
        else:
            files_layout.add_widget(Label(
                text=f"No PowerPoint files found in:\n{folder_path}",
                font_size=self.get_adaptive_font_size(14),
                color=(1, 1, 1, 1),
                size_hint=(1, None),
                height=dp(80),
                halign='center',
                valign='middle',
                text_size=(Window.width * 0.8, None)
            ))

        scroll_view.add_widget(files_layout)
        popup_content.add_widget(scroll_view)

        close_button = Button(
            text="❌ Close",
            size_hint=(1, None),
            height=dp(50),
            font_size=self.get_adaptive_font_size(14)
        )
        self._style_pill_button(close_button, (0.8, 0.2, 0.2, 1), (1, 1, 1, 1))

        # Responsive popup
        popup = Popup(
            title=title,
            content=popup_content,
            size_hint=(0.95, 0.9),
            title_size=self.get_adaptive_font_size(14)
        )
        close_button.bind(on_release=popup.dismiss)
        popup_content.add_widget(close_button)
        popup.open()

    def _handle_file_selected(self, file_name, quarter_num=None, subject=None, week_num=None):
        BASE_DIR = "GregorELibrary"

        file_path = os.path.join(
            BASE_DIR,
            f"Quarter {quarter_num}",
            f"Subject {subject}",
            f"Week {week_num}",
            file_name
        )

        if os.path.exists(file_path):
            try:
                self.loading_popup = self._show_loading_popup()
                if file_name.lower().endswith('.pptx'):
                    Thread(target=self.convert_and_show_ppt, args=(file_path,)).start()
                elif file_name.lower().endswith('.pdf'):
                    Thread(target=self.convert_and_show_pdf, args=(file_path,)).start()
            except Exception as e:
                self._show_error_popup("Error", f"Failed to process file: {e}")
                if self.loading_popup:
                    self.loading_popup.dismiss()
        else:
            self._show_error_popup("Error", f"File not found: {file_path}")

    def pdf_to_images(self, pdf_path):
        """Convert PDF pages to images using PyMuPDF."""
        try:
            doc = fitz.open(pdf_path)
            image_paths = []
            temp_dir = tempfile.mkdtemp()

            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                # Increase resolution with zoom matrix
                zoom_matrix = fitz.Matrix(2.0, 2.0)  # 200% zoom for better quality
                pix = page.get_pixmap(matrix=zoom_matrix)
                img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                temp_file = os.path.join(temp_dir, f"page_{page_num + 1}.png")
                img.save(temp_file)
                image_paths.append(temp_file)
            return image_paths
        except Exception as e:
            print(f"Error rendering PDF: {e}")
            return []

    def convert_and_show_pdf(self, pdf_path):
        """Handle PDF conversion and display."""
        try:
            image_paths = self.pdf_to_images(pdf_path)
            if image_paths:
                Clock.schedule_once(lambda dt: self.show_image_scrollview(image_paths))
            else:
                Clock.schedule_once(
                    lambda dt: self._show_error_popup("Error", "Failed to render PDF"))
        except Exception as e:
            Clock.schedule_once(lambda dt: self._show_error_popup("Error", str(e)))
        finally:
            Clock.schedule_once(
                lambda dt: self.loading_popup.dismiss() if self.loading_popup else None)

    def show_image_scrollview(self, image_paths):  # Renamed from show_ppt_scrollview
        scroll_view = ScrollView(size_hint=(1, 1), do_scroll_x=False, do_scroll_y=True)
        content_layout = BoxLayout(orientation='vertical', size_hint_y=None, spacing=dp(15), padding=dp(10))
        content_layout.bind(minimum_height=content_layout.setter('height'))

        screen_width = Window.width * 0.9
        image_height = dp(int(screen_width * 0.75))  # 4:3 aspect ratio

        for path in image_paths:
            image = AsyncImage(
                source=path,
                allow_stretch=True,
                size_hint=(1, None),
                height=image_height
            )
            content_layout.add_widget(image)

        content_layout.height = len(image_paths) * (image_height + dp(15)) + dp(20)
        scroll_view.add_widget(content_layout)

        close_button = Button(
            text='❌ Close',
            size_hint=(1, None),
            height=dp(50),
            font_size=self.get_adaptive_font_size(14)
        )
        self._style_pill_button(close_button, (0.8, 0.2, 0.2, 1), (1, 1, 1, 1))

        main_layout = BoxLayout(orientation='vertical', spacing=dp(10))
        main_layout.add_widget(scroll_view)
        main_layout.add_widget(close_button)

        popup = Popup(
            title=self._get_responsive_title("Document Viewer"),
            content=main_layout,
            size_hint=(0.95, 0.95),
            title_size=self.get_adaptive_font_size(14)
        )
        close_button.bind(on_release=popup.dismiss)
        popup.open()

    def _show_loading_popup(self):
        popup_content = BoxLayout(orientation='vertical', padding=dp(20), spacing=dp(15))
        self._set_dark_bg(popup_content)
        popup_content.add_widget(Label(
            text="Loading PPT...\nPlease wait",
            font_size=self.get_adaptive_font_size(14),
            color=(1, 1, 1, 1),
            halign='center'
        ))
        # Responsive popup
        popup = Popup(
            title="Loading",
            content=popup_content,
            size_hint=(0.7, 0.3),
            title_size=self.get_adaptive_font_size(14)
        )
        popup.open()
        return popup

    def ppt_to_images(self, ppt_path):
        """Render PPTX slides as images using python-pptx and Pillow."""
        try:
            # Try using Win32COM for better rendering if available
            try:
                # Only try Win32COM on Windows
                if os.name == 'nt':
                    import win32com.client
                    import pythoncom

                    # Initialize COM in the thread
                    pythoncom.CoInitialize()

                    # Get absolute path
                    abs_path = os.path.abspath(ppt_path)

                    # Create temporary directory for images
                    temp_dir = tempfile.mkdtemp()
                    image_paths = []

                    try:
                        # Use PowerPoint COM object for better rendering
                        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
                        powerpoint.Visible = False

                        presentation = powerpoint.Presentations.Open(abs_path)

                        # Save each slide as an image
                        for i in range(1, presentation.Slides.Count + 1):
                            temp_file = os.path.join(temp_dir, f"slide_{i}.png")
                            presentation.Slides(i).Export(temp_file, "PNG", 1024, 768)
                            image_paths.append(temp_file)

                        presentation.Close()
                        powerpoint.Quit()

                        return image_paths

                    except Exception as e:
                        print(f"COM error: {e}")
                        # Fall back to basic rendering if COM fails
                        return self._basic_ppt_rendering(ppt_path)
                    finally:
                        # Clean up COM
                        try:
                            pythoncom.CoUninitialize()
                        except:
                            pass
                else:
                    # On non-Windows platforms, use basic rendering
                    return self._basic_ppt_rendering(ppt_path)

            except ImportError:
                print("Win32com not available, using basic rendering")
                return self._basic_ppt_rendering(ppt_path)
        except Exception as e:
            print(f"Error rendering PPTX: {e}")
            return []

    def _basic_ppt_rendering(self, ppt_path):
        """Basic fallback method using python-pptx and PIL."""
        try:
            presentation = Presentation(ppt_path)
            image_paths = []
            temp_dir = tempfile.mkdtemp()

            # Determine optimal slide dimensions based on screen size
            screen_width = Window.width * 0.9  # 90% of screen width
            slide_width = int(presentation.slide_width / 9525)
            slide_height = int(presentation.slide_height / 9525)

            # Calculate aspect ratio for proper scaling
            aspect_ratio = slide_height / slide_width
            new_width = min(int(screen_width), slide_width)
            new_height = int(new_width * aspect_ratio)

            for i, slide in enumerate(presentation.slides):
                slide_image = Image.new("RGB", (slide_width, slide_height), (255, 255, 255))
                draw = ImageDraw.Draw(slide_image)

                for shape in slide.shapes:
                    # Convert position and size to pixels
                    left = int(shape.left / 9525)
                    top = int(shape.top / 9525)
                    width = int(shape.width / 9525)
                    height = int(shape.height / 9525)

                    # Handle text frames
                    if shape.has_text_frame:
                        text = shape.text_frame.text
                        try:
                            # For Android compatibility, use default font
                            font = ImageFont.load_default()
                            font_size = 12
                        except IOError:
                            font = ImageFont.load_default()
                            font_size = 12

                        # Draw text with better line breaks
                        lines = []
                        current_line = ""
                        words = text.split()
                        for word in words:
                            test_line = current_line + " " + word if current_line else word
                            # Use a more basic approach for text width
                            if len(test_line) * (font_size * 0.6) < width:
                                current_line = test_line
                            else:
                                lines.append(current_line)
                                current_line = word
                        if current_line:
                            lines.append(current_line)

                        # Draw each line
                        line_height = font_size + 4
                        for j, line in enumerate(lines):
                            draw.text((left, top + j * line_height), line, fill="black")

                    # Handle images (basic approach)
                    elif hasattr(shape, "image") and hasattr(shape.image, "blob"):
                        try:
                            img_data = shape.image.blob
                            img = Image.open(io.BytesIO(img_data))
                            img = img.resize((width, height))
                            slide_image.paste(img, (left, top))
                        except Exception as e:
                            print(f"Error with image in slide {i + 1}: {e}")

                # Add slide number
                draw.text((10, slide_height - 20), f"Slide {i + 1}", fill="black")

                # Save slide image
                temp_file = os.path.join(temp_dir, f"slide_{i + 1}.png")
                # Resize to fit screen if needed
                if slide_width > new_width:
                    slide_image = slide_image.resize((new_width, new_height))
                slide_image.save(temp_file)
                image_paths.append(temp_file)

            return image_paths
        except Exception as e:
            print(f"Error in basic rendering: {e}")
            return []

    # Add this helper method to the MainMenu class
    def _get_responsive_title(self, title):
        """Create a responsive title that fits the window width by truncating if necessary."""
        # Get available width (considering padding in the popup)
        available_width = Window.width * 0.85  # 85% of window width to account for margins

        # Estimate characters that can fit based on average character width
        # Using a conservative estimate based on font size
        font_size = self.get_adaptive_font_size(14)
        chars_per_width = int(available_width / (font_size * 0.6))  # Approximate character width

        # If title is too long, truncate it intelligently
        if len(title) > chars_per_width:
            # Split the title by parts
            title_parts = title.split(" - ")

            # Keep the first and last parts, abbreviate middle parts if needed
            if len(title_parts) > 2:
                # Always show the current level (last part)
                last_part = title_parts[-1]
                first_part = title_parts[0]

                # Calculate remaining space for middle parts
                remaining_chars = chars_per_width - len(first_part) - len(last_part) - 7  # 7 chars for " - ... - "

                if remaining_chars > 10:  # If we have reasonable space
                    middle_parts = title_parts[1:-1]
                    middle_str = ""

                    # Try to fit as many middle parts as possible
                    for part in middle_parts:
                        if len(middle_str) + len(part) + 3 <= remaining_chars:  # +3 for " - "
                            middle_str += " - " + part if middle_str else part
                        else:
                            # Add ellipsis if we can't fit all parts
                            middle_str += "..." if not middle_str else "..."
                            break

                    return f"{first_part} - {middle_str} - {last_part}"
                else:
                    # Very limited space, just show first and last with ellipsis
                    return f"{first_part} - ... - {last_part}"
            else:
                # Only two parts, try to show both
                first_part = title_parts[0]
                last_part = title_parts[1]

                if len(first_part) + len(last_part) + 3 > chars_per_width:
                    # Still too long, truncate both parts
                    chars_each = (chars_per_width - 3) // 2  # -3 for " - "
                    return f"{first_part[:chars_each]}... - ...{last_part[-chars_each:]}"

        # Title fits, return as is
        return title

    # Now update the popup creation in all relevant methods:

    def _show_quarters_popup(self, title):
        self.current_navigation_path = title

        popup_content = BoxLayout(orientation='vertical', padding=dp(20), spacing=dp(15))
        self._set_dark_bg(popup_content)

        # Use responsive title
        responsive_title = self._get_responsive_title(f"{title} - Select Quarter")

        popup_content.add_widget(Label(
            text=responsive_title,
            font_size=self.get_adaptive_font_size(14),
            color=(1, 1, 1, 1),
            size_hint=(1, None),
            height=dp(40)
        ))

        # Rest of the method remains the same
        quarters = ["1st", "2nd", "3rd", "4th"]
        for i, quarter in enumerate(quarters, start=1):
            quarter_button = Button(
                text=f"📖 {quarter} Quarter",
                size_hint=(1, None),
                height=dp(50),
                font_size=self.get_adaptive_font_size(14)
            )
            quarter_num = i
            quarter_button.bind(on_release=lambda btn, q=quarter_num:
            self._show_subjects_popup(f"{title} - {quarters[q - 1]} Quarter", q))
            self._style_pill_button(quarter_button, (0.5, 0.5, 0.5, 1), (1, 1, 1, 1))
            popup_content.add_widget(quarter_button)

        close_button = Button(
            text="❌ Close",
            size_hint=(1, None),
            height=dp(50),
            font_size=self.get_adaptive_font_size(14)
        )
        self._style_pill_button(close_button, (0.8, 0.2, 0.2, 1), (1, 1, 1, 1))

        # Use responsive title for popup
        popup = Popup(
            title=self._get_responsive_title(title),
            content=popup_content,
            size_hint=(0.95, 0.8),
            title_size=self.get_adaptive_font_size(14)
        )
        close_button.bind(on_release=popup.dismiss)
        popup_content.add_widget(close_button)
        popup.open()

    def _show_subjects_popup(self, title, quarter_num=None):
        # Store current navigation path for later use
        self.current_navigation_path = title

        # Extract quarter from title if not provided
        if quarter_num is None:
            quarter_part = title.split(" - ")[1].split()[0][0]
            quarter_num = int(quarter_part)

        popup_content = BoxLayout(orientation='vertical', padding=dp(20), spacing=dp(15))
        self._set_dark_bg(popup_content)

        # Use responsive title
        responsive_title = self._get_responsive_title(f"{title} - Select Subject")

        popup_content.add_widget(Label(
            text=responsive_title,
            font_size=self.get_adaptive_font_size(14),
            color=(1, 1, 1, 1),
            size_hint=(1, None),
            height=dp(40)
        ))

        if quarter_num in [1, 2]:  # 1st and 2nd quarter subjects
            subjects = [
                "1 - Oral Comm",
                "2 - Komunikasyon",
                "3 - General Mathematics",
                "4 - Earth Science",
                "5 - Understanding Culture, Society and Politics",
                "6 - Personal Development",
                "7 - PE"
            ]
        elif quarter_num in [3, 4]:  # 3rd and 4th quarter subjects
            subjects = [
                "1 - Reading and Writing",
                "2 - Pagbasa",
                "3 - 21st Century Literature",
                "4 - Stats and Probability",
                "5 - PE"
            ]

        scroll_height = min(Window.height * 0.5, dp(300))
        scroll_view = ScrollView(size_hint=(1, None), height=scroll_height)
        subjects_layout = BoxLayout(orientation='vertical', size_hint_y=None, spacing=dp(15))
        subjects_layout.bind(minimum_height=subjects_layout.setter('height'))

        for subject in subjects:
            subject_button = Button(
                text=f"📚 Subject {subject}",
                size_hint=(1, None),
                height=dp(50),
                font_size=self.get_adaptive_font_size(14)
            )
            subject_button.bind(on_release=lambda btn, q=quarter_num, s=subject:
            self._show_weeks_popup(f"{title} - Subject {s}", q, s))
            self._style_pill_button(subject_button, (0.6, 0.6, 0.6, 1), (1, 1, 1, 1))
            subjects_layout.add_widget(subject_button)

        scroll_view.add_widget(subjects_layout)
        popup_content.add_widget(scroll_view)

        close_button = Button(
            text="❌ Close",
            size_hint=(1, None),
            height=dp(50),
            font_size=self.get_adaptive_font_size(14)
        )
        self._style_pill_button(close_button, (0.8, 0.2, 0.2, 1), (1, 1, 1, 1))

        # Use responsive title for popup
        popup = Popup(
            title=self._get_responsive_title(title),
            content=popup_content,
            size_hint=(0.95, 0.9),
            title_size=self.get_adaptive_font_size(14)
        )
        close_button.bind(on_release=popup.dismiss)
        popup_content.add_widget(close_button)
        popup.open()

    def _show_weeks_popup(self, title, quarter_num=None, subject=None):
        self.current_navigation_path = title

        # Extract quarter and subject if not provided
        if quarter_num is None:
            title_parts = title.split(" - ")
            quarter_part = title_parts[1].split()[0][0]
            quarter_num = int(quarter_part)

        if subject is None:
            subject = title.split(" - ")[-1]

        popup_content = BoxLayout(orientation='vertical', padding=dp(20), spacing=dp(15))
        self._set_dark_bg(popup_content)

        # Use responsive title
        responsive_title = self._get_responsive_title(f"{title} - Select Week")

        popup_content.add_widget(Label(
            text=responsive_title,
            font_size=self.get_adaptive_font_size(14),
            color=(1, 1, 1, 1),
            size_hint=(1, None),
            height=dp(40)
        ))

        # Rest of the method remains the same
        scroll_height = min(Window.height * 0.6, dp(400))
        scroll_view = ScrollView(size_hint=(1, None), height=scroll_height)
        weeks_layout = BoxLayout(orientation='vertical', size_hint_y=None, spacing=dp(15))
        weeks_layout.bind(minimum_height=weeks_layout.setter('height'))

        weeks = [f"Week {i}" for i in range(1, 9)]
        for week in weeks:
            week_button = Button(
                text=f"📅 {week}",
                size_hint=(1, None),
                height=dp(50),
                font_size=self.get_adaptive_font_size(14)
            )
            week_num = int(week.split()[1])
            week_button.bind(on_release=lambda btn, q=quarter_num, s=subject, w=week_num:
            self._show_files_popup(f"{title} - Week {w}", q, s, w))
            self._style_pill_button(week_button, (0.7, 0.7, 0.7, 1), (1, 1, 1, 1))
            weeks_layout.add_widget(week_button)

        scroll_view.add_widget(weeks_layout)
        popup_content.add_widget(scroll_view)

        close_button = Button(
            text="❌ Close",
            size_hint=(1, None),
            height=dp(50),
            font_size=self.get_adaptive_font_size(14)
        )
        self._style_pill_button(close_button, (0.8, 0.2, 0.2, 1), (1, 1, 1, 1))

        # Use responsive title for popup
        popup = Popup(
            title=self._get_responsive_title(title),
            content=popup_content,
            size_hint=(0.95, 0.9),
            title_size=self.get_adaptive_font_size(14)
        )
        close_button.bind(on_release=popup.dismiss)
        popup_content.add_widget(close_button)
        popup.open()

    def _show_files_popup(self, title, quarter_num=None, subject=None, week_num=None):
        self.current_navigation_path = title

        # Extract information if not provided
        if any(param is None for param in [quarter_num, subject, week_num]):
            title_parts = title.split(" - ")
            if quarter_num is None:
                quarter_part = title_parts[1].split()[0][0]
                quarter_num = int(quarter_part)

            if subject is None:
                subject = title_parts[2].replace("Subject ", "")

            if week_num is None:
                week_num = int(title_parts[3].split()[1])

        popup_content = BoxLayout(orientation='vertical', padding=dp(20), spacing=dp(15))
        self._set_dark_bg(popup_content)

        # Use responsive title
        responsive_title = self._get_responsive_title(f"{title} - Select File")

        popup_content.add_widget(Label(
            text=responsive_title,
            font_size=self.get_adaptive_font_size(14),
            color=(1, 1, 1, 1),
            size_hint=(1, None),
            height=dp(40)
        ))

        # Rest of the method remains the same
        scroll_height = min(Window.height * 0.6, dp(400))
        scroll_view = ScrollView(size_hint=(1, None), height=scroll_height)
        files_layout = BoxLayout(orientation='vertical', size_hint_y=None, spacing=dp(15))
        files_layout.bind(minimum_height=files_layout.setter('height'))

        BASE_DIR = "GregorELibrary"

        folder_path = os.path.join(
            BASE_DIR,
            f"Quarter {quarter_num}",
            f"Subject {subject}",
            f"Week {week_num}"
        )

        print(f"Looking for files in: {folder_path}")
        print(f"Quarter: {quarter_num}, Subject: {subject}, Week: {week_num}")

        files = []
        if os.path.exists(folder_path):
            files = [f for f in os.listdir(folder_path) if f.endswith('.pptx')]

        if files:
            for file in files:
                file_button = Button(
                    text=f"📄 {file}",
                    size_hint=(1, None),
                    height=dp(50),
                    font_size=self.get_adaptive_font_size(14)
                )
                file_button.bind(on_release=lambda btn, f=file, q=quarter_num, s=subject, w=week_num:
                self._handle_file_selected(f, q, s, w))
                self._style_pill_button(file_button, (0.8, 0.8, 0.8, 1), (1, 1, 1, 1))
                files_layout.add_widget(file_button)
        else:
            files_layout.add_widget(Label(
                text=f"No PowerPoint files found in:\n{folder_path}",
                font_size=self.get_adaptive_font_size(14),
                color=(1, 1, 1, 1),
                size_hint=(1, None),
                height=dp(80),
                halign='center',
                valign='middle',
                text_size=(Window.width * 0.8, None)
            ))

        scroll_view.add_widget(files_layout)
        popup_content.add_widget(scroll_view)

        close_button = Button(
            text="❌ Close",
            size_hint=(1, None),
            height=dp(50),
            font_size=self.get_adaptive_font_size(14)
        )
        self._style_pill_button(close_button, (0.8, 0.2, 0.2, 1), (1, 1, 1, 1))

        # Use responsive title for popup
        popup = Popup(
            title=self._get_responsive_title(title),
            content=popup_content,
            size_hint=(0.95, 0.9),
            title_size=self.get_adaptive_font_size(14)
        )
        close_button.bind(on_release=popup.dismiss)
        popup_content.add_widget(close_button)
        popup.open()

    # Make sure to also update the show_ppt_scrollview method
    def show_ppt_scrollview(self, image_paths):
        scroll_view = ScrollView(size_hint=(1, 1), do_scroll_x=False, do_scroll_y=True)
        content_layout = BoxLayout(orientation='vertical', size_hint_y=None, spacing=dp(15), padding=dp(10))
        content_layout.bind(minimum_height=content_layout.setter('height'))

        # Calculate adaptive height based on screen width and aspect ratio
        screen_width = Window.width * 0.9
        image_height = dp(int(screen_width * 0.75))  # 4:3 aspect ratio is common for slides

        for path in image_paths:
            image = AsyncImage(
                source=path,
                allow_stretch=True,
                size_hint=(1, None),
                height=image_height
            )
            content_layout.add_widget(image)

        # Set total height
        content_layout.height = len(image_paths) * (image_height + dp(15)) + dp(20)
        scroll_view.add_widget(content_layout)

        close_button = Button(
            text='❌ Close',
            size_hint=(1, None),
            height=dp(50),
            font_size=self.get_adaptive_font_size(14)
        )
        self._style_pill_button(close_button, (0.8, 0.2, 0.2, 1), (1, 1, 1, 1))

        main_layout = BoxLayout(orientation='vertical', spacing=dp(10))
        main_layout.add_widget(scroll_view)
        main_layout.add_widget(close_button)

        # Use responsive title for the PPT viewer popup
        popup = Popup(
            title=self._get_responsive_title("PPT Viewer"),
            content=main_layout,
            size_hint=(0.95, 0.95),
            title_size=self.get_adaptive_font_size(14)
        )
        close_button.bind(on_release=popup.dismiss)
        popup.open()

    def convert_and_show_ppt(self, ppt_path):
        try:
            image_paths = self.ppt_to_images(ppt_path)
            if image_paths:
                Clock.schedule_once(lambda dt: self.show_ppt_scrollview(image_paths))
            else:
                Clock.schedule_once(lambda dt: self._show_error_popup("Error", "Failed to render PPTX"))
        except Exception as e:
            Clock.schedule_once(lambda dt: self._show_error_popup("Error", str(e)))
        finally:
            Clock.schedule_once(lambda dt: self.loading_popup.dismiss() if self.loading_popup else None)

    def _show_error_popup(self, title, message):
        popup_content = BoxLayout(orientation='vertical', padding=dp(20), spacing=dp(10))
        self._set_dark_bg(popup_content)

        popup_content.add_widget(Label(
            text=message,
            font_size=sp(14),
            color=(1, 1, 1, 1)
        ))

        close_button = Button(
            text="Close",
            size_hint=(1, None),
            height=dp(40),
            font_size=sp(14)
        )
        self._style_pill_button(close_button, (0.8, 0.2, 0.2, 1), (1, 1, 1, 1))

        popup = Popup(title=title, content=popup_content, size_hint=(0.8, 0.4))
        close_button.bind(on_release=popup.dismiss)
        popup_content.add_widget(close_button)
        popup.open()

    def show_about_us(self, instance):
        popup_content = BoxLayout(orientation='vertical', padding=dp(20), spacing=dp(10))
        self._set_dark_bg(popup_content)

        popup_content.add_widget(Label(
            text="About Us",
            font_size=sp(24),
            color=(1, 1, 1, 1)
        ))
        popup_content.add_widget(Label(
            text="We are a team dedicated to providing easy access\n"
                    "to educational resources. Our platform is\n"
                    "designed to help students manage study materials\n"
                    "anytime, anywhere.",
            font_size=sp(14),
            color=(1, 1, 1, 1)
        ))

        close_button = Button(
            text="Close",
            size_hint=(1, None),
            height=dp(40),
            font_size=sp(14)
        )
        self._style_pill_button(close_button, (0.8, 0.2, 0.2, 1), (1, 1, 1, 1))

        popup = Popup(title="About Us", content=popup_content, size_hint=(0.9, 0.9))
        close_button.bind(on_release=popup.dismiss)
        popup_content.add_widget(close_button)
        popup.open()

    def _set_dark_bg(self, layout):
        with layout.canvas.before:
            Color(0, 0, 0, 1)
            layout.bg_rect = Rectangle(size=layout.size, pos=layout.pos)
        layout.bind(size=self._update_bg_rect_popup, pos=self._update_bg_rect_popup)

    def _update_bg_rect_popup(self, layout, *args):
        layout.bg_rect.pos = layout.pos
        layout.bg_rect.size = layout.size

    def _style_pill_button(self, btn, bg_color, text_color):
        btn.background_color = (0, 0, 0, 0)
        btn.color = text_color
        btn.background_normal = ''
        btn.background_down = ''

        with btn.canvas.before:
            Color(*bg_color)
            btn.rect = RoundedRectangle(size=btn.size, pos=btn.pos, radius=[dp(20)])
        btn.bind(size=self._update_btn_rect, pos=self._update_btn_rect)

    def _update_btn_rect(self, instance, *args):
        instance.rect.pos = instance.pos
        instance.rect.size = instance.size


class GregorELibraryApp(App):
    def build(self):
        # Load initial configuration
        Window.bind(on_keyboard=self.on_keyboard)
        return MainMenu()

    def on_keyboard(self, window, key, *args):
        # Disable back button exit
        if key == 27:
            return True
        return False

if __name__ == "__main__":
    GregorELibraryApp().run()
